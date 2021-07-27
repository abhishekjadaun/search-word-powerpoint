from tkinter import *
from pptx import Presentation
import os
from tkinter import filedialog
# filename=""
def UploadAction(event=None):
    global filename
    filename = filedialog.askopenfilename()
    

def search_text():
    findSlideId=[]
    # filename = filedialog.askopenfilename()
    path=filename
    keyword=e2.get().lower()
    path = filename.replace('/','\\\\')
    
    prs = Presentation(path) 
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = shape.text.lower()
                if keyword in shape.text:
                    # print(files)
                    slide_id = prs.slides.index(slide)
                    findSlideId.append(slide_id+1)
                    break
    if not findSlideId:
        values = 'No Record Match'
    else:
        values = '\n'.join(str(v) for v in findSlideId)
    label_text.set(values)
    
 
window = Tk()
window.wm_title("Word Search in PowerPoint")
window.minsize(800,700)
label_text=StringVar()
Label(window, text="Select File:").grid(row=0, sticky=W)
Label(window, text="enter Word:").grid(row=1, sticky=W)
Label(window, text="Result of Slide No:").grid(row=3, sticky=W)
result=Label(window, text="", textvariable=label_text).grid(row=3,column=1, sticky=W)
button = Button(window, text='upload', command=UploadAction)
button.grid(row=0, column=2,columnspan=2)

e2 = Entry(window,width = 100)
e2.grid(row=1, column=1, columnspan=30)
 
b = Button(window, text="Search", command=search_text)
b.grid(row=2, column=2,columnspan=2)
 
 
window.mainloop()

